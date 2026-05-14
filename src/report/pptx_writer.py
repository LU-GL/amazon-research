"""PPT报告生成模块。

用python-pptx生成市调分析PPT报告。
"""
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData

DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
ACCENT_BLUE = RGBColor(0x4F, 0x81, 0xBD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)


def write_market_research_ppt(
    matrix_df: pd.DataFrame,
    sentiment_data: dict[str, dict],
    comparison: dict | None = None,
    category: str = "",
    output_path: Path | None = None,
) -> Path:
    """生成市场调研PPT报告。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    today = datetime.now().strftime("%Y-%m-%d")

    # Slide 1: 封面
    _add_title_slide(prs, category, today)

    # Slide 2: 市场概览
    _add_overview_slide(prs, matrix_df, category)

    # Slide 3: 竞品矩阵表
    _add_matrix_slide(prs, matrix_df)

    # Slide 4: 痛点分析
    _add_pain_points_slide(prs, sentiment_data, comparison)

    # Slide 5: 好评分析
    _add_praise_slide(prs, sentiment_data)

    # Slide 6: 功能需求 & 机会点
    _add_opportunities_slide(prs, sentiment_data, comparison)

    # Slide 7: 战略建议
    _add_recommendations_slide(prs, comparison)

    # Slide 8: 附录
    _add_appendix_slide(prs, today)

    # 保存
    if output_path is None:
        from src.config import OUTPUT_DIR
        output_path = OUTPUT_DIR / "market_research.pptx"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"PPT已保存: {output_path}")
    return output_path


def _add_title_slide(prs, category: str, date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景色
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    # 标题
    title_text = f"{category + ' ' if category else ''}市场调研报告"
    _add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(2),
                  title_text, size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 日期
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1),
                  f"生成日期: {date}", size=18, color=WHITE, alignment=PP_ALIGN.CENTER)


def _add_overview_slide(prs, matrix_df: pd.DataFrame, category: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "市场概览")

    total = len(matrix_df)
    prices = pd.to_numeric(matrix_df["价格"].str.replace(r"[^\d.]", "", regex=True), errors="coerce").dropna()
    ratings = pd.to_numeric(matrix_df["评分"], errors="coerce").dropna()

    lines = [
        f"分析产品数: {total}",
        f"品类: {category or '未指定'}",
    ]
    if not prices.empty:
        lines.append(f"价格范围: ${prices.min():.2f} ~ ${prices.max():.2f}")
        lines.append(f"平均价格: ${prices.mean():.2f}")
    if not ratings.empty:
        lines.append(f"评分范围: {ratings.min():.1f} ~ {ratings.max():.1f}")
        lines.append(f"平均评分: {ratings.mean():.2f}")

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  "\n".join(lines), size=20, color=BLACK)


def _add_matrix_slide(prs, matrix_df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "竞品矩阵对比")

    cols = ["ASIN", "品牌", "价格", "评分", "评论数", "核心痛点1", "好评主题1"]
    available = [c for c in cols if c in matrix_df.columns]
    df_sub = matrix_df[available].head(10)

    rows_count = len(df_sub) + 1
    cols_count = len(available)
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(12)
    height = Inches(0.4 * rows_count)

    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table

    # 表头
    for j, col in enumerate(available):
        cell = table.cell(0, j)
        cell.text = col
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # 数据行
    for i, (_, row) in enumerate(df_sub.iterrows()):
        for j, col in enumerate(available):
            cell = table.cell(i + 1, j)
            val = str(row.get(col, ""))
            if len(val) > 30:
                val = val[:30] + "..."
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)


def _add_pain_points_slide(prs, sentiment_data: dict, comparison: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "痛点分析")

    lines = []
    if comparison and comparison.get("common_pain_points"):
        lines.append("=== 跨产品共同痛点 ===")
        for pp in comparison["common_pain_points"][:5]:
            affected = len(pp.get("asins", []))
            lines.append(f"  {pp['theme']} (影响{affected}个产品, 频次{pp['frequency']})")

    if comparison and comparison.get("unique_pain_points"):
        lines.append("\n=== 差异化痛点 ===")
        for pp in comparison["unique_pain_points"][:5]:
            lines.append(f"  {pp['theme']} (仅{pp['asins'][0] if pp.get('asins') else '?'}), 频次{pp['frequency']})")

    if not lines:
        lines.append("暂无跨产品痛点对比数据")

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  "\n".join(lines), size=16, color=BLACK)


def _add_praise_slide(prs, sentiment_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "好评分析")

    all_praise = {}
    for asin, analysis in sentiment_data.items():
        for pp in analysis.get("praise_points", []):
            theme = pp["theme"]
            if theme not in all_praise:
                all_praise[theme] = {"frequency": 0, "asins": []}
            all_praise[theme]["frequency"] += pp.get("frequency", 0)
            all_praise[theme]["asins"].append(asin)

    sorted_praise = sorted(all_praise.items(), key=lambda x: x[1]["frequency"], reverse=True)

    lines = []
    for theme, info in sorted_praise[:8]:
        lines.append(f"  {theme} (频次{info['frequency']}, {len(info['asins'])}个产品)")

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  "\n".join(lines) if lines else "暂无好评数据", size=16, color=BLACK)


def _add_opportunities_slide(prs, sentiment_data: dict, comparison: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "市场机会点")

    lines = []
    if comparison and comparison.get("feature_opportunities"):
        lines.append("=== 用户功能需求 ===")
        for fr in comparison["feature_opportunities"][:5]:
            lines.append(f"  {fr['theme']} (频次{fr['frequency']})")

    # 质量问题
    quality_issues = []
    for asin, analysis in sentiment_data.items():
        for qi in analysis.get("quality_issues", []):
            quality_issues.append(qi)
    if quality_issues:
        lines.append("\n=== 质量改进机会 ===")
        sorted_qi = sorted(quality_issues, key=lambda x: x.get("frequency", 0), reverse=True)
        seen = set()
        for qi in sorted_qi[:5]:
            theme = qi["theme"]
            if theme not in seen:
                seen.add(theme)
                lines.append(f"  {theme} (严重度: {qi.get('severity', '?')})")

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  "\n".join(lines) if lines else "暂无机会点数据", size=16, color=BLACK)


def _add_recommendations_slide(prs, comparison: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "战略建议")

    lines = []
    if comparison and comparison.get("insights"):
        for i, insight in enumerate(comparison["insights"], 1):
            lines.append(f"{i}. {insight}")
    else:
        lines.append("基于以上分析数据，建议关注以下方向：")
        lines.append("1. 解决市场共同痛点以获取竞争优势")
        lines.append("2. 关注高频功能需求，打造差异化产品")
        lines.append("3. 在竞品薄弱环节重点突破")

    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  "\n".join(lines), size=16, color=BLACK)


def _add_appendix_slide(prs, date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, "附录")
    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(5),
                  f"报告生成日期: {date}\n"
                  f"数据来源: 卖家精灵导出 + Amazon Listing抓取\n"
                  f"分析工具: Claude API 情感分析\n"
                  f"分析维度: 痛点/好评/功能需求/质量改进",
                  size=14, color=BLACK)


def _add_slide_title(slide, text: str):
    """添加幻灯片标题栏。"""
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), slide.slide_width, Inches(1.2)  # Rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.8)
    tf.vertical_anchor = 1  # middle


def _add_text_box(slide, left, top, width, height, text, size=14,
                  bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        p.space_after = Pt(4)
